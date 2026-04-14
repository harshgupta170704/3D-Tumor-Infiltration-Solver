from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
import re

def create_presentation(md_path, output_path):
    # Initialize presentation (widescreen 16:9)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    with open(md_path, 'r', encoding='utf-8') as f:
        content = f.read()

    # Split by Slide headers
    slides_content = re.split(r'## Slide \d+:', content)
    
    # We ignore the very first part (since it's just theme suggestions)
    if len(slides_content) > 1:
        slides_content = slides_content[1:]

    # Track available images to inject
    images = {
        "segmentation": "final_output/segmentation.png",
        "tumor_density": "final_output/tumor_density.png",
        "physics": "final_output/physics_params.png",
        "growth": "final_output/growth_simulation.png"
    }

    for idx, slide_text in enumerate(slides_content):
        # Extract title
        lines = slide_text.strip().split('\n')
        title_line = [l for l in lines if l.startswith('*   **Title:**')]
        title = lines[0].strip()
        if title_line:
            title = title_line[0].replace('*   **Title:**', '').strip()
            
        # For Slide 1, use Title Layout
        if idx == 0:
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = title
            subtitle = slide.placeholders[1]
            sub_text = [l for l in lines if l.startswith('*   **Subtitle:**')]
            if sub_text:
                subtitle.text = sub_text[0].replace('*   **Subtitle:**', '').strip()
            continue

        # Standard Text Layout
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # Add basic title
        if slide.shapes.title:
            slide.shapes.title.text = title
            
        # Parse content
        content_box = slide.placeholders[1]
        tf = content_box.text_frame
        tf.clear()
        
        in_content = False
        notes = []
        for line in lines:
            if line.startswith('*   **Content:**'):
                in_content = True
                continue
            if line.startswith('*   **Visual:**') or line.startswith('*   **Speaker Notes:**'):
                in_content = False
                if line.startswith('*   **Speaker Notes:**'):
                    notes.append(line.replace('*   **Speaker Notes:**', '').strip())
            
            if in_content and line.strip().startswith('*'):
                p = tf.add_paragraph()
                p.text = line.replace('**', '').replace('*', '').strip()
                p.level = 0
            elif in_content and line.strip() != "":
                p = tf.add_paragraph()
                p.text = line.replace('**', '').strip()
                p.level = 1

        # Add speaker notes
        if slide.has_notes_slide and notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = " ".join(notes)

        # Inject real images based on slide context
        added_img = False
        img_to_add = None
        
        if "Ground Truth Labels" in title or "Segmentation" in title:
            img_to_add = images["segmentation"]
        elif "MRI Looks Like" in title or "Density" in title:
            img_to_add = images["tumor_density"]
        elif "Physics Model" in title or "Parameters" in title:
            img_to_add = images["physics"]
        elif "Output" in title or "Results" in title or "Trajectory" in title:
            img_to_add = images["growth"]
            
        if img_to_add and os.path.exists(img_to_add):
            try:
                # Add image to right side
                left = Inches(6.5)
                top = Inches(2.0)
                height = Inches(4.5)
                slide.shapes.add_picture(img_to_add, left, top, height=height)
                added_img = True
            except Exception as e:
                print(f"Could not add image {img_to_add} to slide {idx}: {e}")

    # Save
    prs.save(output_path)
    print(f"Successfully generated {output_path}")

if __name__ == '__main__':
    create_presentation('presentation_content.md', 'Final_Project_Presentation.pptx')
